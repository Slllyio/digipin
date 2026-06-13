# -*- coding: utf-8 -*-
"""
DigiPin Urban Intelligence - Sales Deck Generator
Exhaustive feature-per-slide PPT positioned as premium product.
"""

import sys
if sys.platform == "win32":
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ═══ THEME COLORS ═══
CYAN = RGBColor(0x00, 0xE5, 0xFF)
PURPLE = RGBColor(0x7C, 0x4D, 0xFF)
DARK_BG = RGBColor(0x0A, 0x0A, 0x1A)
CARD_BG = RGBColor(0x14, 0x14, 0x2A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xE0, 0xE0, 0xE0)
MUTED = RGBColor(0x99, 0x99, 0x99)
DIM = RGBColor(0x55, 0x55, 0x55)
ACCENT_GREEN = RGBColor(0x00, 0xC8, 0x53)
ACCENT_ORANGE = RGBColor(0xFF, 0xAB, 0x40)
ACCENT_RED = RGBColor(0xFF, 0x52, 0x52)

TOTAL_SLIDES = 24
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# ═══ SCREENSHOT MAPPING (slide_num -> image_path) ═══
# Fresh captures in ppt-screenshots/ + existing project root screenshots
SCREENSHOTS = {
    1:  "ppt-screenshots/01-homepage.png",       # Title — full homepage
    3:  "ppt-screenshots/01-homepage.png",       # What We Built — overview
    4:  "digipin-zoomed.png",                    # Grid System — zoomed grid cells
    5:  "digipin-dashboard.png",                 # City Coverage — dashboard view
    6:  "digipin-cell-clicked.png",              # Smart Search — cell selected
    7:  "ppt-screenshots/02-detail-panel.png",   # Intelligence Profile — detail panel
    8:  "ppt-screenshots/03-scores.png",         # Quality Scores — scores + radar
    9:  "ppt-screenshots/03-scores.png",         # Radar Chart — same view
    10: "google_buildings_flat.png",             # Building Intelligence — flat view
    11: "ppt-screenshots/04-buildings-3d.png",   # 3D Visualization — 3D buildings
    12: "roads-active.png",                      # Road Network — color-coded roads
    13: "lulc-overlay-test.png",                 # Climate & Land Use — LULC overlay
    14: "layers-all-active.png",                 # Digital Twin — multiple layers active
    15: "multi-layers-active.png",               # 3D Heatmap — multi-layer view
    16: "digipin-v3-sectors.png",                # Urban Query Engine — sectors panel
    17: "digipin-landing-analysis.png",          # Location Compare — analysis view
    18: "landuse-boundaries-weather.png",        # Walkability Isochrone — boundaries
    19: "floating-dialogs-test.png",             # Bookmarks & Reports — floating dialogs
    20: "digipin-expanded-final.png",            # DISHA AI — expanded panel
    22: "digipin-v2-final.png",                  # Use Cases — full platform view
    24: "ppt-screenshots/01-homepage.png",       # Closing — homepage
}


# ═══════════════════════════════════════════════════════════
# HELPERS
# ═══════════════════════════════════════════════════════════

def bg(slide):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = DARK_BG

def txt(slide, l, t, w, h, text, sz=14, c=WHITE, b=False, a=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = c
    p.font.bold = b
    p.font.name = "Calibri"
    p.alignment = a
    return tb

def multi_txt(slide, l, t, w, h, lines):
    """Add text box with multiple styled lines. Each line: (text, size, color, bold)"""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (text, sz, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(4)
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = text
        r.font.size = Pt(sz)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.name = "Calibri"
    return tb

def card(slide, l, t, w, h, fill=CARD_BG):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    s.shadow.inherit = False
    return s

def stat(slide, l, t, val, lbl, vc=CYAN, w=Inches(2.2), h=Inches(1.6)):
    card(slide, l, t, w, h)
    txt(slide, l, t + Inches(0.2), w, Inches(0.7), val, sz=38, c=vc, b=True, a=PP_ALIGN.CENTER)
    txt(slide, l, t + Inches(0.95), w, Inches(0.4), lbl.upper(), sz=10, c=MUTED, a=PP_ALIGN.CENTER)

def snum(slide, n):
    txt(slide, Inches(14.2), Inches(0.3), Inches(1.5), Inches(0.3),
        f"{n:02d} / {TOTAL_SLIDES}", sz=11, c=DIM, a=PP_ALIGN.RIGHT)

def tag(slide, text, l=Inches(5.5), t=Inches(2.0)):
    txt(slide, l, t, Inches(5), Inches(0.35), text.upper(), sz=11, c=CYAN, b=True, a=PP_ALIGN.CENTER)

def title(slide, text, t=Inches(2.5)):
    txt(slide, Inches(1.5), t, Inches(13), Inches(1.0), text, sz=44, c=WHITE, b=True, a=PP_ALIGN.CENTER)

def subtitle(slide, text, t=Inches(3.7)):
    txt(slide, Inches(2.5), t, Inches(11), Inches(0.8), text, sz=18, c=MUTED, a=PP_ALIGN.CENTER)

def bullets(slide, l, t, items, w=Inches(6.5), sz=15):
    tb = slide.shapes.add_textbox(l, t, w, Inches(len(items) * 0.48))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        p.space_before = Pt(4)
        r = p.add_run()
        r.text = item
        r.font.size = Pt(sz)
        r.font.color.rgb = LIGHT
        r.font.name = "Calibri"

def scorebar(slide, l, t, label, val, w=Inches(5.5)):
    bh = Inches(0.12)
    txt(slide, l, t, Inches(1.8), Inches(0.3), label, sz=13, c=MUTED, a=PP_ALIGN.RIGHT)
    track = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l+Inches(2.0), t+Inches(0.08), w-Inches(2.7), bh)
    track.fill.solid(); track.fill.fore_color.rgb = RGBColor(0x1A,0x1A,0x2E); track.line.fill.background()
    bw = int((w - Inches(2.7)) * (val / 100))
    if bw > 0:
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l+Inches(2.0), t+Inches(0.08), bw, bh)
        bar.fill.solid(); bar.fill.fore_color.rgb = CYAN; bar.line.fill.background()
    txt(slide, l+w-Inches(0.5), t, Inches(0.5), Inches(0.3), str(val), sz=13, c=CYAN, b=True)

def fcard(slide, l, t, ic, ttl, desc, w=Inches(4.6), h=Inches(2.5)):
    card(slide, l, t, w, h)
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, l+Inches(0.3), t+Inches(0.25), Inches(0.55), Inches(0.55))
    circ.fill.solid(); circ.fill.fore_color.rgb = RGBColor(0x1A,0x10,0x3A); circ.line.fill.background()
    tf = circ.text_frame; tf.paragraphs[0].alignment = PP_ALIGN.CENTER; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    r = tf.paragraphs[0].add_run(); r.text = ic; r.font.size = Pt(16); r.font.color.rgb = CYAN
    txt(slide, l+Inches(0.3), t+Inches(1.0), w-Inches(0.6), Inches(0.35), ttl, sz=17, c=WHITE, b=True)
    txt(slide, l+Inches(0.3), t+Inches(1.4), w-Inches(0.6), Inches(0.9), desc, sz=12, c=MUTED)

def arch_box(slide, l, t, ttl, sub, w=Inches(2.4)):
    card(slide, l, t, w, Inches(1.1))
    txt(slide, l, t+Inches(0.2), w, Inches(0.35), ttl, sz=14, c=WHITE, b=True, a=PP_ALIGN.CENTER)
    txt(slide, l, t+Inches(0.55), w, Inches(0.3), sub, sz=10, c=MUTED, a=PP_ALIGN.CENTER)

def img(slide, slide_num, l=Inches(8.5), t=Inches(1.5), w=Inches(6.8), h=Inches(5.5)):
    """Add screenshot image to slide if available. Returns True if added."""
    path = SCREENSHOTS.get(slide_num)
    if path and os.path.exists(path):
        slide.shapes.add_picture(path, l, t, w, h)
        return True
    return False

def img_full(slide, slide_num, l=Inches(0), t=Inches(0), w=Inches(16), h=Inches(9)):
    """Add full-bleed background screenshot (dimmed by overlay card on top)."""
    path = SCREENSHOTS.get(slide_num)
    if path and os.path.exists(path):
        slide.shapes.add_picture(path, l, t, w, h)
        return True
    return False

def img_card(slide, slide_num, l=Inches(8.3), t=Inches(1.3), w=Inches(7.2), h=Inches(6.0)):
    """Add screenshot inside a card frame."""
    card(slide, l, t, w, h)
    path = SCREENSHOTS.get(slide_num)
    if path and os.path.exists(path):
        # Inset the image slightly within the card
        slide.shapes.add_picture(path, l+Inches(0.15), t+Inches(0.15),
                                  w-Inches(0.3), h-Inches(0.3))
        return True
    return False


# ═══════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); snum(s, 1)
tag(s, "Location Intelligence Platform", t=Inches(1.6))
txt(s, Inches(1), Inches(2.2), Inches(14), Inches(1.3),
    "DigiPin", sz=80, c=CYAN, b=True, a=PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(3.4), Inches(14), Inches(0.8),
    "Urban Intelligence", sz=48, c=WHITE, b=True, a=PP_ALIGN.CENTER)
txt(s, Inches(2.5), Inches(4.5), Inches(11), Inches(0.8),
    "The most comprehensive micro-location intelligence engine for Indian cities.\n"
    "160+ data dimensions. 30+ quality scores. One platform.",
    sz=19, c=MUTED, a=PP_ALIGN.CENTER)
for i, (v, l) in enumerate([("160+", "Data Points / Cell"), ("30+", "Intelligence Scores"),
                              ("528K", "Building Footprints"), ("12", "Metro Cities")]):
    stat(s, Inches(2.4 + i*2.8), Inches(6.0), v, l)

# ═══════════════════════════════════════════════════════════
# SLIDE 2: THE GAP
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); snum(s, 2)
tag(s, "The Market Gap")
title(s, "Location decisions are still\nmade on gut feeling")
subtitle(s, "India's $200B real estate market and 4,000+ urban local bodies make critical decisions\n"
            "without granular, multi-dimensional location data. The data exists -- but it's scattered\n"
            "across dozens of sources, requires GIS expertise, and costs months to assemble.")
for i, (v, l) in enumerate([("$200B", "Real Estate Market"), ("4,000+", "Urban Local Bodies"),
                              ("Months", "To Assemble Data"), ("$50K+", "Per City Analysis")]):
    stat(s, Inches(2.4 + i*2.8), Inches(6.0), v, l, vc=ACCENT_RED)

# ═══════════════════════════════════════════════════════════
# SLIDE 3: WHAT IS DIGIPIN
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); snum(s, 3)
tag(s, "What We Built")
title(s, "Micro-location intelligence,\nautomated at scale")
txt(s, Inches(2), Inches(4.0), Inches(12), Inches(1.2),
    "DigiPin fuses satellite imagery, building footprints, road networks, climate classifications,\n"
    "land use data, and points of interest into a unified intelligence layer.\n\n"
    "Every 100m x 100m cell in the city gets a comprehensive quality profile --\n"
    "computed automatically, updated continuously, queryable instantly.",
    sz=17, c=LIGHT, a=PP_ALIGN.CENTER)
for i, (v, l) in enumerate([("5+", "Data Pipelines"), ("160+", "Computed Features"),
                              ("Sub-second", "Query Response"), ("100m", "Grid Resolution")]):
    stat(s, Inches(2.4 + i*2.8), Inches(6.2), v, l)

# ═══════════════════════════════════════════════════════════
# SLIDE 4: DIGIPIN GRID SYSTEM
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); snum(s, 4)
tag(s, "Proprietary Grid", l=Inches(0.8), t=Inches(1.0))
txt(s, Inches(0.8), Inches(1.6), Inches(7), Inches(0.7),
    "The DigiPin Encoding System", sz=36, c=WHITE, b=True)
txt(s, Inches(0.8), Inches(2.5), Inches(7), Inches(1.5),
    "Built on India Post's DigiPin standard, our grid uses a proprietary 16-symbol alphabet "
    "(2-9, C, F, J, K, L, M, P, T) with 4x4 recursive subdivision.\n\n"
    "Each 10-character code maps to a unique ~100m cell anywhere in India. "
    "This gives us 1 billion+ addressable micro-locations with deterministic, "
    "offline-capable encoding -- no database lookup required.",
    sz=15, c=LIGHT)
bullets(s, Inches(0.8), Inches(4.8), [
    "16-symbol alphabet -- compact, unambiguous, human-readable",
    "4x4 recursive grid -- each level divides into 16 sub-cells",
    "10-character precision -- ~100m x 100m resolution",
    "Deterministic encoding -- works offline, no API dependency",
    "Covers all of India -- lat 2.5-38.5, lon 63.5-99.5",
], w=Inches(7))

# Right side - screenshot of zoomed grid
img_card(s, 4)

# ═══════════════════════════════════════════════════════════
# SLIDE 5: CITY COVERAGE
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); snum(s, 5)
img(s, 5, l=Inches(9), t=Inches(0), w=Inches(7), h=Inches(9))  # Right half background
tag(s, "Coverage")
title(s, "12 Metro Cities, Full Coverage")
subtitle(s, "Every grid cell in each city is pre-computed with complete feature and score profiles.\n"
            "Instant city switching with zero load time -- all data is pre-indexed and cache-optimized.")
cities = ["Indore", "Bengaluru", "Mumbai", "Delhi NCR", "Hyderabad", "Chennai",
          "Pune", "Ahmedabad", "Jaipur", "Kolkata", "Lucknow", "Bhopal"]
for i, city in enumerate(cities):
    col, row = i % 6, i // 6
    card(s, Inches(1.0 + col*2.4), Inches(5.5 + row*1.6), Inches(2.1), Inches(1.3))
    txt(s, Inches(1.0 + col*2.4), Inches(5.8 + row*1.6), Inches(2.1), Inches(0.5),
        city, sz=16, c=WHITE, b=True, a=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SLIDE 6: SMART SEARCH
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); snum(s, 6)
tag(s, "Feature Deep-Dive", l=Inches(0.8), t=Inches(1.0))
txt(s, Inches(0.8), Inches(1.6), Inches(7), Inches(0.7),
    "Intelligent Search", sz=38, c=WHITE, b=True)
txt(s, Inches(0.8), Inches(2.5), Inches(7), Inches(1.0),
    "Two search modes that cover every use case: natural language place search "
    "with geocoding, and precision DigiPin code lookup for exact cell targeting.\n\n"
    "The map smoothly flies to the result with cinematic animation, "
    "and the grid cell is auto-selected for immediate profiling.",
    sz=15, c=LIGHT)

card(s, Inches(0.8), Inches(4.5), Inches(6.5), Inches(3.5))
multi_txt(s, Inches(1.3), Inches(4.8), Inches(5.5), Inches(3.0), [
    ("Place Name Search", 16, CYAN, True),
    ("Type any landmark, address, or area name.", 13, MUTED, False),
    ("Geocoded in real-time with auto-complete.", 13, MUTED, False),
    ("", 8, MUTED, False),
    ("DigiPin Code Search", 16, CYAN, True),
    ("Enter a 10-character DigiPin code.", 13, MUTED, False),
    ("Instant cell lookup -- no network call required.", 13, MUTED, False),
    ("Shareable: send a code, share an exact location.", 13, MUTED, False),
])

# Right - screenshot of cell selected after search
img_card(s, 6)

# ═══════════════════════════════════════════════════════════
# SLIDE 7: CELL INTELLIGENCE PROFILE
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); snum(s, 7)
tag(s, "Core Product")
title(s, "The Intelligence Profile", t=Inches(1.0))
txt(s, Inches(2), Inches(2.0), Inches(12), Inches(0.8),
    "Click any cell on the map. In under a second, DigiPin computes and displays\n"
    "a comprehensive intelligence dossier with 160+ features and 30+ quality scores.",
    sz=17, c=MUTED, a=PP_ALIGN.CENTER)

profile_items = [
    ("30+", "Quality Scores", "Livability, walkability, safety,\ngreen cover, noise, connectivity"),
    ("160+", "Raw Features", "Buildings, roads, amenities,\ntransit, terrain, land use"),
]
for i, (v, t, d) in enumerate(profile_items):
    l = Inches(0.5 + i*3.9)
    card(s, l, Inches(3.3), Inches(3.5), Inches(4.5))
    txt(s, l, Inches(3.7), Inches(3.5), Inches(0.7), v, sz=42, c=CYAN, b=True, a=PP_ALIGN.CENTER)
    txt(s, l, Inches(4.5), Inches(3.5), Inches(0.4), t, sz=17, c=WHITE, b=True, a=PP_ALIGN.CENTER)
    txt(s, l+Inches(0.3), Inches(5.2), Inches(2.9), Inches(1.5), d, sz=14, c=MUTED, a=PP_ALIGN.CENTER)

# Right - screenshot of detail panel
img_card(s, 7, l=Inches(8.3), t=Inches(1.0), w=Inches(7.2), h=Inches(7.2))

# ═══════════════════════════════════════════════════════════
# SLIDE 8: QUALITY SCORES DETAIL
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); snum(s, 8)
tag(s, "Scoring Engine", l=Inches(0.8), t=Inches(0.8))
txt(s, Inches(0.8), Inches(1.4), Inches(7), Inches(0.7),
    "30+ Calibrated Quality Scores", sz=36, c=WHITE, b=True)
txt(s, Inches(0.8), Inches(2.3), Inches(7), Inches(1.0),
    "Each score is computed from multiple underlying features using weighted formulas "
    "calibrated against ground truth. Scores range from 0 to 100 and are comparable "
    "across cells, neighborhoods, and cities.\n\n"
    "The scoring engine processes satellite-derived classifications, "
    "building morphology data, road network topology, and point-of-interest density "
    "to produce actionable intelligence.",
    sz=14, c=LIGHT)

categories = [
    "Livability -- holistic quality of life metric",
    "Walkability -- pedestrian infrastructure & connectivity",
    "Green Cover Index -- vegetation density & park access",
    "Safety Score -- lighting, CCTV, emergency services proximity",
    "Commercial Viability -- footfall, competition, amenity mix",
    "Connectivity -- road density, transit access, highway proximity",
    "Infrastructure -- utilities, civic amenities, development",
    "Noise Exposure -- traffic, commercial, industrial noise sources",
]
bullets(s, Inches(0.8), Inches(4.2), categories, w=Inches(7), sz=14)

# Right - screenshot of actual scores with radar chart
img_card(s, 8)

# ═══════════════════════════════════════════════════════════
# SLIDE 9: RADAR CHART
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); snum(s, 9)
tag(s, "Visual Intelligence", l=Inches(0.8), t=Inches(1.0))
txt(s, Inches(0.8), Inches(1.6), Inches(7), Inches(0.7),
    "The Location Fingerprint", sz=38, c=WHITE, b=True)
txt(s, Inches(0.8), Inches(2.5), Inches(7), Inches(1.5),
    "Every location has a unique quality fingerprint -- a radar chart that shows "
    "all dimensions simultaneously. This makes it trivially easy to:\n\n"
    "- Spot strengths and weaknesses at a glance\n"
    "- Compare two locations visually\n"
    "- Communicate quality profiles to non-technical stakeholders\n"
    "- Identify areas that need specific interventions\n\n"
    "The radar chart is interactive -- hover any axis for the detailed score breakdown "
    "and the raw features that contribute to it.",
    sz=15, c=LIGHT)

# Right - screenshot of actual radar chart + scores dialog
img_card(s, 9)

# ═══════════════════════════════════════════════════════════
# SLIDE 10: BUILDING INTELLIGENCE
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); snum(s, 10)
tag(s, "Building Analytics")
title(s, "528,000 Building Footprints\nWith Structural Intelligence")
txt(s, Inches(2), Inches(4.0), Inches(12), Inches(0.8),
    "Every building in the coverage area has been individually cataloged with footprint geometry,\n"
    "estimated height, structural classification, and Local Climate Zone membership.",
    sz=17, c=MUTED, a=PP_ALIGN.CENTER)

for i, (v, l, d) in enumerate([
    ("528K", "Footprints", "Individual building\ngeometries cataloged"),
    ("Height", "Estimation", "Per-building height\nfrom ML models"),
]):
    l2 = Inches(0.5 + i*3.9)
    card(s, l2, Inches(5.2), Inches(3.5), Inches(3.0))
    txt(s, l2, Inches(5.5), Inches(3.5), Inches(0.6), v, sz=36, c=CYAN, b=True, a=PP_ALIGN.CENTER)
    txt(s, l2, Inches(6.2), Inches(3.5), Inches(0.4), l, sz=15, c=WHITE, b=True, a=PP_ALIGN.CENTER)
    txt(s, l2+Inches(0.3), Inches(6.7), Inches(2.9), Inches(1.0), d, sz=13, c=MUTED, a=PP_ALIGN.CENTER)

# Right - screenshot of building footprints
img_card(s, 10, l=Inches(8.3), t=Inches(1.0), w=Inches(7.2), h=Inches(7.2))

# ═══════════════════════════════════════════════════════════
# SLIDE 11: 3D VISUALIZATION
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); snum(s, 11)
tag(s, "3D Engine", l=Inches(0.8), t=Inches(1.0))
txt(s, Inches(0.8), Inches(1.6), Inches(7), Inches(0.7),
    "3D City Visualization", sz=38, c=WHITE, b=True)
txt(s, Inches(0.8), Inches(2.5), Inches(7), Inches(2.0),
    "Every building is extruded to its actual estimated height and color-coded "
    "by structural type. The 3D engine supports:\n\n"
    "- Real-time pitch and rotation with cinematic camera controls\n"
    "- Building type classification (residential, commercial, industrial)\n"
    "- Height-based color gradient for urban density analysis\n"
    "- Seamless toggle between 2D footprint and 3D extrusion views\n"
    "- Full integration with all overlay layers\n\n"
    "This gives stakeholders a photorealistic understanding of urban morphology "
    "without ever visiting the site.",
    sz=15, c=LIGHT)

# Right - screenshot of 3D buildings extruded
img_card(s, 11)

# ═══════════════════════════════════════════════════════════
# SLIDE 12: ROAD NETWORK
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); snum(s, 12)
tag(s, "Transport Layer")
title(s, "Complete Road Network Analysis")
txt(s, Inches(2), Inches(3.5), Inches(12), Inches(0.8),
    "Every road segment is classified by type and rendered with distinct color coding.\n"
    "Road density, connectivity, and accessibility metrics feed directly into quality scores.",
    sz=17, c=MUTED, a=PP_ALIGN.CENTER)

roads = [
    ("Motorways", ACCENT_RED),
    ("Primary", ACCENT_ORANGE),
    ("Secondary", RGBColor(0xFF,0xD7,0x40)),
    ("Residential", ACCENT_GREEN),
]
for i, (name, color) in enumerate(roads):
    l = Inches(0.5 + i*1.9)
    card(s, l, Inches(4.8), Inches(1.6), Inches(1.0))
    card(s, l+Inches(0.2), Inches(5.0), Inches(1.2), Inches(0.15), fill=color)
    txt(s, l, Inches(5.3), Inches(1.6), Inches(0.4), name, sz=14, c=WHITE, b=True, a=PP_ALIGN.CENTER)

# Right - screenshot of color-coded road network
img_card(s, 12, l=Inches(8.3), t=Inches(1.0), w=Inches(7.2), h=Inches(7.2))

# ═══════════════════════════════════════════════════════════
# SLIDE 13: CLIMATE & LAND USE OVERLAYS
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); snum(s, 13)
tag(s, "Geospatial Overlays")
title(s, "Satellite-Derived Intelligence Layers")
txt(s, Inches(2), Inches(3.5), Inches(12), Inches(0.8),
    "Three critical overlay layers transform raw satellite data into actionable urban intelligence.\n"
    "Each layer required months of processing, calibration, and validation.",
    sz=17, c=MUTED, a=PP_ALIGN.CENTER)

overlay_names = [
    ("Local Climate Zones", "17 classes"),
    ("Land Use / Land Cover", "54 classes"),
    ("Ward Boundaries", "Admin zones"),
]
for i, (name, sub) in enumerate(overlay_names):
    l = Inches(0.5 + i*2.5)
    card(s, l, Inches(4.8), Inches(2.2), Inches(1.2))
    txt(s, l, Inches(5.0), Inches(2.2), Inches(0.4), name, sz=14, c=CYAN, b=True, a=PP_ALIGN.CENTER)
    txt(s, l, Inches(5.4), Inches(2.2), Inches(0.3), sub, sz=12, c=MUTED, a=PP_ALIGN.CENTER)

# Right - screenshot of LULC overlay
img_card(s, 13, l=Inches(8.3), t=Inches(1.0), w=Inches(7.2), h=Inches(7.2))

# ═══════════════════════════════════════════════════════════
# SLIDE 14: DIGITAL TWIN LAYERS
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); snum(s, 14)
tag(s, "Digital Twin", l=Inches(0.8), t=Inches(1.0))
txt(s, Inches(0.8), Inches(1.6), Inches(7), Inches(0.7),
    "Multi-Source Data Fusion", sz=38, c=WHITE, b=True)
txt(s, Inches(0.8), Inches(2.5), Inches(7), Inches(1.5),
    "DigiPin integrates data from the Overture Maps Foundation -- "
    "a collaboration between Amazon, Meta, Microsoft, and TomTom -- "
    "to provide the most comprehensive geospatial layer stack available.\n\n"
    "Each layer can be toggled independently, and all layers are "
    "pre-indexed for instant rendering at any zoom level.",
    sz=15, c=LIGHT)

# Right - screenshot of all layers active
img_card(s, 14)

# ═══════════════════════════════════════════════════════════
# SLIDE 15: 3D HEATMAP ANALYSIS
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); snum(s, 15)
tag(s, "Spatial Analytics")
title(s, "3D Heatmap Engine")
txt(s, Inches(2), Inches(3.5), Inches(12), Inches(0.8),
    "Select any quality metric and instantly visualize its spatial distribution across the entire city.\n"
    "3D columns rise from the map -- taller columns represent higher scores. The worst zones are immediately obvious.",
    sz=17, c=MUTED, a=PP_ALIGN.CENTER)

metrics = ["Livability", "Walkability", "Green Cover", "Safety",
           "Commercial", "Connectivity", "Infrastructure", "Healthcare"]
for i, m in enumerate(metrics):
    col, row = i % 4, i // 4
    card(s, Inches(0.5 + col*1.9), Inches(5.0 + row*1.5), Inches(1.7), Inches(1.2))
    txt(s, Inches(0.5 + col*1.9), Inches(5.3 + row*1.5), Inches(1.7), Inches(0.5),
        m, sz=13, c=WHITE, b=True, a=PP_ALIGN.CENTER)

# Right - screenshot of multi-layer heatmap view
img_card(s, 15, l=Inches(8.3), t=Inches(1.0), w=Inches(7.2), h=Inches(7.2))

# ═══════════════════════════════════════════════════════════
# SLIDE 16: URBAN QUERY ENGINE
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); snum(s, 16)
tag(s, "Decision Engine", l=Inches(0.8), t=Inches(0.8))
txt(s, Inches(0.8), Inches(1.4), Inches(7), Inches(0.7),
    "52 Urban Analytical Queries", sz=36, c=WHITE, b=True)
txt(s, Inches(0.8), Inches(2.3), Inches(7), Inches(1.2),
    "The query engine answers the hardest location questions automatically. "
    "Select a query, and the engine samples 25 points across the visible area, "
    "computes 160+ features per point, applies weighted scoring, "
    "and ranks the top locations -- all in under 30 seconds.\n\n"
    "This analysis would take a human analyst 2-3 weeks per query.",
    sz=15, c=LIGHT)

sectors = [
    ("Commercial", "8"),
    ("Residential", "6"),
    ("Infrastructure", "6"),
    ("Environment", "6"),
    ("Real Estate", "8"),
    ("Tourism", "5"),
    ("Social", "6"),
]
for i, (name, count) in enumerate(sectors):
    col, row = i % 4, i // 4
    l = Inches(0.8 + col * 1.8)
    t = Inches(4.5 + row * 1.5)
    card(s, l, t, Inches(1.6), Inches(1.2))
    txt(s, l, t+Inches(0.2), Inches(1.6), Inches(0.35), name, sz=13, c=CYAN, b=True, a=PP_ALIGN.CENTER)
    txt(s, l, t+Inches(0.6), Inches(1.6), Inches(0.3), f"{count} queries", sz=11, c=MUTED, a=PP_ALIGN.CENTER)

# Right - screenshot of query engine sectors
img_card(s, 16, l=Inches(8.3), t=Inches(1.0), w=Inches(7.2), h=Inches(7.2))

# ═══════════════════════════════════════════════════════════
# SLIDE 17: LOCATION COMPARE
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); snum(s, 17)
tag(s, "Comparison Engine")
title(s, "Side-by-Side Location Intelligence")
txt(s, Inches(2), Inches(3.5), Inches(12), Inches(0.8),
    "Pin up to 3 locations and compare them across every quality dimension.\n"
    "Radar charts, score tables, and feature breakdowns make site selection decisions data-driven.",
    sz=17, c=MUTED, a=PP_ALIGN.CENTER)

# Right - screenshot of analysis/comparison view
img_card(s, 17, l=Inches(8.3), t=Inches(1.0), w=Inches(7.2), h=Inches(7.2))

# Left - comparison summary cards
for i, (name, grade) in enumerate([("Location A", "B+"), ("Location B", "A-"), ("Location C", "B")]):
    t = Inches(4.8 + i*1.3)
    card(s, Inches(0.5), t, Inches(7.4), Inches(1.0))
    txt(s, Inches(1.0), t+Inches(0.15), Inches(2.5), Inches(0.35), name, sz=16, c=WHITE, b=True)
    txt(s, Inches(3.5), t+Inches(0.15), Inches(2.5), Inches(0.35), f"Grade: {grade}", sz=18, c=CYAN, b=True)
    txt(s, Inches(6.0), t+Inches(0.15), Inches(1.5), Inches(0.35), "Compare", sz=12, c=MUTED)

# ═══════════════════════════════════════════════════════════
# SLIDE 18: WALKABILITY ISOCHRONE
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); snum(s, 18)
tag(s, "Accessibility Analysis", l=Inches(0.8), t=Inches(1.0))
txt(s, Inches(0.8), Inches(1.6), Inches(7), Inches(0.7),
    "Walkability Isochrone", sz=38, c=WHITE, b=True)
txt(s, Inches(0.8), Inches(2.5), Inches(7), Inches(1.5),
    "For any selected cell, DigiPin computes walking-distance isochrones "
    "showing exactly how far a person can walk in 5, 10, and 15 minutes.\n\n"
    "This uses real road network topology -- not straight-line distance -- "
    "accounting for actual walkable paths, intersections, and barriers.\n\n"
    "Critical for real estate valuation, retail site selection, "
    "transit planning, and accessibility compliance.",
    sz=15, c=LIGHT)

zones = [
    ("5 min", "400m", ACCENT_GREEN),
    ("10 min", "800m", ACCENT_ORANGE),
    ("15 min", "1.2km", ACCENT_RED),
]
for i, (time, radius, color) in enumerate(zones):
    t = Inches(5.0 + i*1.2)
    card(s, Inches(0.8), t, Inches(6.5), Inches(1.0))
    card(s, Inches(1.1), t+Inches(0.2), Inches(0.15), Inches(0.6), fill=color)
    txt(s, Inches(1.5), t+Inches(0.2), Inches(2.0), Inches(0.4), time, sz=22, c=WHITE, b=True)
    txt(s, Inches(3.8), t+Inches(0.2), Inches(2.0), Inches(0.4), radius, sz=16, c=CYAN)

# Right - screenshot of boundaries/isochrone view
img_card(s, 18)

# ═══════════════════════════════════════════════════════════
# SLIDE 19: BOOKMARKS & REPORTS
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); snum(s, 19)
tag(s, "Productivity Tools")
title(s, "Bookmark, Report, Share")
txt(s, Inches(2), Inches(3.5), Inches(12), Inches(0.5),
    "Save locations, generate professional reports, and build a library of analyzed sites.",
    sz=17, c=MUTED, a=PP_ALIGN.CENTER)

tools = [
    ("Bookmarks", "Save locations with notes. Persists across sessions."),
    ("PDF Reports", "One-click print-ready intelligence reports."),
    ("Share via DigiPin", "Share a 10-char code for exact location."),
    ("Export Data", "Export scores and features to Excel or BI tools."),
]
for i, (name, desc) in enumerate(tools):
    t = Inches(4.5 + i*1.1)
    card(s, Inches(0.5), t, Inches(7.4), Inches(0.9))
    txt(s, Inches(1.0), t+Inches(0.15), Inches(3.0), Inches(0.35), name, sz=17, c=CYAN, b=True)
    txt(s, Inches(1.0), t+Inches(0.5), Inches(6.0), Inches(0.35), desc, sz=13, c=LIGHT)

# Right - screenshot of floating dialogs / bookmarks
img_card(s, 19)

# ═══════════════════════════════════════════════════════════
# SLIDE 20: DISHA AI
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); snum(s, 20)
tag(s, "AI-Powered", l=Inches(0.8), t=Inches(0.8))
txt(s, Inches(0.8), Inches(1.4), Inches(7), Inches(0.7),
    "DISHA: AI Location Analyst", sz=38, c=WHITE, b=True)
txt(s, Inches(0.8), Inches(2.3), Inches(7), Inches(1.5),
    "DISHA is an AI assistant that turns complex location data into natural language insights. "
    "Ask any question about a location, and DISHA analyzes all 160+ features to give a "
    "contextual, data-backed answer.\n\n"
    "Unlike generic chatbots, DISHA has the full intelligence profile injected as context -- "
    "it doesn't hallucinate; it reasons over real data.",
    sz=15, c=LIGHT)

features = [
    "Natural language Q&A about any location",
    "Full 160+ feature context per query",
    "Multi-provider: local LLM or cloud inference",
    "Smart context filtering -- sends only relevant data",
    "Response caching with TTL for performance",
    "City-wide batch scan for pattern detection",
    "Privacy-first: local inference keeps data on-device",
]
bullets(s, Inches(0.8), Inches(4.5), features, w=Inches(7), sz=14)

# Right - screenshot of DISHA / expanded panel
img_card(s, 20)

# ═══════════════════════════════════════════════════════════
# SLIDE 21: TECH ARCHITECTURE
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); snum(s, 21)
tag(s, "Technology")
title(s, "Enterprise-Grade Architecture")
txt(s, Inches(2), Inches(3.3), Inches(12), Inches(0.5),
    "Zero-server deployment. Sub-second response. Works offline. Scales to any city.",
    sz=18, c=MUTED, a=PP_ALIGN.CENTER)

txt(s, Inches(1.5), Inches(4.2), Inches(13), Inches(0.35),
    "RENDERING & VISUALIZATION", sz=11, c=CYAN, b=True, a=PP_ALIGN.CENTER)
for i, (t2, sub) in enumerate([
    ("MapLibre GL JS", "WebGL vector rendering"),
    ("3D Extrusion Engine", "Building height visualization"),
    ("PMTiles", "Serverless tile delivery"),
    ("Real-time Overlays", "Dynamic layer compositing"),
]):
    arch_box(s, Inches(1.2 + i*3.5), Inches(4.7), t2, sub)

txt(s, Inches(1.5), Inches(6.2), Inches(13), Inches(0.35),
    "INTELLIGENCE & COMPUTE", sz=11, c=PURPLE, b=True, a=PP_ALIGN.CENTER)
for i, (t2, sub) in enumerate([
    ("Feature Engine", "160+ computed features"),
    ("Scoring Pipeline", "Weighted quality metrics"),
    ("Query Processor", "25-point spatial sampling"),
    ("DISHA AI", "LLM with data context"),
]):
    arch_box(s, Inches(1.2 + i*3.5), Inches(6.7), t2, sub)

# ═══════════════════════════════════════════════════════════
# SLIDE 22: USE CASES
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); snum(s, 22)
tag(s, "Applications")
title(s, "Who Needs DigiPin?", t=Inches(1.0))

use_cases = [
    ("RE", "Real Estate Developers", "Site selection, land valuation, neighborhood quality scoring for investment decisions worth crores."),
    ("UP", "Urban Local Bodies", "Infrastructure gap analysis, ward-level planning, green cover monitoring, walkability improvement."),
    ("BIZ", "Retail & F&B Chains", "Optimal store placement using 52 analytical queries. Reduce site selection time from months to minutes."),
    ("FIN", "Banks & NBFCs", "Property collateral assessment, area risk scoring, portfolio concentration analysis by location quality."),
    ("INS", "Insurance Companies", "Risk zonation, flood/heat exposure, building density assessment for premium calibration."),
    ("CON", "Consulting Firms", "White-label location intelligence for client advisory. Integrate via API for custom reports."),
]
for i, (ic, name, desc) in enumerate(use_cases):
    col, row = i % 3, i // 3
    fcard(s, Inches(0.5 + col*5.1), Inches(2.5 + row*3.0), ic, name, desc)

# ═══════════════════════════════════════════════════════════
# SLIDE 23: COMPETITIVE ADVANTAGE
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); snum(s, 23)
tag(s, "Why DigiPin", l=Inches(0.8), t=Inches(0.8))
txt(s, Inches(0.8), Inches(1.4), Inches(8), Inches(0.7),
    "What Makes This Hard to Replicate", sz=36, c=WHITE, b=True)

advantages = [
    ("Multi-source data fusion", "Integrating 5+ heterogeneous data pipelines -- satellite imagery, "
     "building footprints, road topology, climate classifications, and POI data -- into a single "
     "coherent intelligence layer is an engineering challenge that takes years to solve."),
    ("Proprietary scoring engine", "The 30+ quality scores aren't simple averages. Each score uses "
     "calibrated weighted formulas with cross-validated feature importance. "
     "The scoring methodology is our core IP."),
    ("Pan-India grid coverage", "Computing 160+ features for every 100m cell across 12 metro cities "
     "requires massive data processing infrastructure. Expanding to new cities takes weeks, not months."),
    ("AI with real context", "DISHA doesn't guess -- it reasons over actual data. "
     "Injecting 160+ structured features as LLM context produces answers that are "
     "verifiable and trustworthy, unlike generic AI recommendations."),
    ("Sub-second performance", "Pre-computed features and scores with aggressive caching "
     "deliver instant results. No waiting for API calls or database queries."),
]
for i, (name, desc) in enumerate(advantages):
    t = Inches(2.4 + i*1.3)
    card(s, Inches(0.8), t, Inches(14.4), Inches(1.1))
    txt(s, Inches(1.2), t+Inches(0.15), Inches(3.5), Inches(0.35), name, sz=16, c=CYAN, b=True)
    txt(s, Inches(1.2), t+Inches(0.5), Inches(13.5), Inches(0.5), desc, sz=13, c=LIGHT)

# ═══════════════════════════════════════════════════════════
# SLIDE 24: CLOSING
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); snum(s, 24)
tag(s, "Let's Talk", t=Inches(2.0))
txt(s, Inches(1), Inches(2.6), Inches(14), Inches(1.2),
    "DigiPin", sz=76, c=CYAN, b=True, a=PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(3.8), Inches(14), Inches(0.8),
    "Urban Intelligence", sz=44, c=WHITE, b=True, a=PP_ALIGN.CENTER)
txt(s, Inches(2.5), Inches(4.8), Inches(11), Inches(0.8),
    "The most comprehensive micro-location intelligence platform for Indian cities.\n"
    "Making every location decision data-driven.",
    sz=19, c=MUTED, a=PP_ALIGN.CENTER)

for i, (v, l) in enumerate([("160+", "Data Dimensions"), ("30+", "Quality Scores"),
                              ("528K", "Buildings Profiled"), ("52", "Ready Queries")]):
    stat(s, Inches(2.4 + i*2.8), Inches(6.0), v, l)

txt(s, Inches(3), Inches(8.2), Inches(10), Inches(0.4),
    "Micro-Location Intelligence  |  Multi-Source Data Fusion  |  AI-Powered Insights",
    sz=13, c=DIM, a=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════
out = "DigiPin-Presentation.pptx"
prs.save(out)
sz = os.path.getsize(out) / 1024
print(f"\n  PowerPoint saved: {out}")
print(f"  Slides: {TOTAL_SLIDES}")
print(f"  Size: {sz:.0f} KB\n")
